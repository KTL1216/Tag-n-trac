from datetime import datetime, timezone
import os
import re
import sys
import json
import numpy as np
import matplotlib.pyplot as plt
import requests
import json
import pandas as pd
import openpyxl
from datetime import timedelta
from pptx import Presentation
from pptx.util import Inches
from openpyxl.utils.dataframe import dataframe_to_rows
import matplotlib.dates as mdates

prs = Presentation()

## login
API_BASE = "https://api.tagntrac.io"

# Placeholder variables for user credentials and filename
fname = ""
id = "username" 
pwd = "password" 
def prompt():
    """Prompt user for username, password, and file name for device id list."""
    id = input("Enter username: ")
    pwd = input("Enter password: ")
    fname = input("Enter IMEI list file (default imei.txt): ")
    return id, pwd, fname

def login2(email, password):
    login_response = requests.post(f"{API_BASE}/login?clientId=Tbocs0cjhrac",
                             data = json.dumps({"emailId" : email, "userSecret" : password,"reqType": "cognitoAuth"}),
                             headers={"Content-Type" : "application/json", "Origin" : "DOC.API"})
    try:
        if login_response.json()["status"] == "SUCCESS":
            print("Login successful as ", email)
            return (login_response.json()["idToken"], login_response.json()['clientApiKey']['clientId'])
    except Exception as e:
        print(f"Exception: {str(e)}")
    print(f"Login failed: {login_response.text}")
    return (None, None)

# Capture user input
id, pwd, fname = prompt()
if fname == "":
    fname = "imei.txt"

idToken, xapikey2 = login2(id, pwd)
common_headers2 = {"Authorization" : idToken,
                  "Origin" : f"{API_BASE}",
                  "x-api-key" : xapikey2}


def generate_time_string(hours_ago):
    # Get the current time in UTC
    end_time = datetime.now(timezone.utc)
    # Calculate the start time by subtracting the given hours
    start_time = end_time - timedelta(hours=hours_ago)

    # Format both times to the ISO 8601 format with milliseconds
    start_str = start_time.strftime('%Y-%m-%dT%H:%M:%S.000Z')
    end_str = end_time.strftime('%Y-%m-%dT%H:%M:%S.000Z')

    # Construct the final string
    result_string = f"?start={start_str}&end={end_str}"
    return result_string


def get_device_data_v2(device_id, hours_ago):
    queryDates = generate_time_string(hours_ago)
    print(queryDates)
    response = requests.get(f"{API_BASE}/v2/device/{device_id}/data"+queryDates,
                            headers=common_headers2)
    if response.json()['status'] == 'SUCCESS':
        data = response.json()['response']
        return data
    else: 
        print("Get Device data2 failed: "+device_id)
        return None

def vbat_data(data_entry, dev_id, hours_ago):
    timestamp_ms = int(data_entry['ts'])
    timestamp_s = timestamp_ms / 1000
    reported_time = datetime.fromtimestamp(timestamp_s, timezone.utc)

    # Calculate time difference
    current_utc_time = datetime.now(timezone.utc)
    time_difference = current_utc_time - reported_time

    if data_entry['vbat'] is not None:
        # Calculate hours, minutes, and seconds from time difference
        hours_formatted = time_difference.total_seconds() // 3600
        minutes = (time_difference.total_seconds() % 3600) // 60
        seconds = time_difference.total_seconds() % 60
        # Format the time difference
        formatted_time_difference = f"{int(hours_formatted)} hrs {int(minutes)} mins {int(seconds)} secs"
        data_dict = {
            'IMEI': dev_id,
            'Timestamp': reported_time.strftime('%Y-%m-%d %H:%M:%S'),
            'Time passed since Reported': formatted_time_difference,
            'vBat': int(data_entry['vbat']),
            f'Reported time since {hours_ago}hrs ago (hrs)': float(hours_ago)-time_difference.total_seconds() / 3600
        }
        return data_dict
    else:
        return None 
    
def rsrp_data(data_entry, dev_id, hours_ago):
    timestamp_ms = int(data_entry['ts'])
    timestamp_s = timestamp_ms / 1000
    reported_time = datetime.fromtimestamp(timestamp_s, timezone.utc)

    # Calculate time difference
    current_utc_time = datetime.now(timezone.utc)
    time_difference = current_utc_time - reported_time

    if data_entry['rsrp'] is not None:
        # Calculate hours, minutes, and seconds from time difference
        hours_formatted = time_difference.total_seconds() // 3600
        minutes = (time_difference.total_seconds() % 3600) // 60
        seconds = time_difference.total_seconds() % 60
        # Format the time difference
        formatted_time_difference = f"{int(hours_formatted)} hrs {int(minutes)} mins {int(seconds)} secs"
        data_dict = {
            'IMEI': dev_id,
            'Timestamp': reported_time.strftime('%Y-%m-%d %H:%M:%S'),
            'Time passed since Reported': formatted_time_difference,
            'RSRP': int(data_entry['rsrp']),
            f'Reported time since {hours_ago}hrs ago (hrs)': float(hours_ago)-time_difference.total_seconds() / 3600
        }
        return data_dict
    else:
        return None 
    
# Convert 'Time passed since Reported' into a total seconds for plotting
def convert_to_seconds(t):
    try:
        time_parts = {'hrs': 0, 'mins': 0, 'secs': 0}
        parts = t.split()
        for i in range(0, len(parts), 2):
            if parts[i + 1].startswith('hr'):
                time_parts['hrs'] = int(parts[i])
            elif parts[i + 1].startswith('min'):
                time_parts['mins'] = int(parts[i])
            elif parts[i + 1].startswith('sec'):
                time_parts['secs'] = int(parts[i])
        return time_parts['hrs'] * 3600 + time_parts['mins'] * 60 + time_parts['secs']
    except Exception as e:
        print(f"Error converting time: {t} - {e}")
        return 0  # return 0 if there's an error, or you could choose to handle it differently

def convert_to_hours(t):
    try:
        time_parts = {'hrs': 0, 'mins': 0, 'secs': 0}
        parts = t.split()
        for i in range(0, len(parts), 2):
            if parts[i + 1].startswith('hr'):
                time_parts['hrs'] = int(parts[i])
            elif parts[i + 1].startswith('min'):
                time_parts['mins'] = int(parts[i])
            elif parts[i + 1].startswith('sec'):
                time_parts['secs'] = int(parts[i])
        total_seconds = time_parts['hrs'] * 3600 + time_parts['mins'] * 60 + time_parts['secs']
        return total_seconds / 3600  # Convert seconds to hours
    except Exception as e:
        print(f"Error converting time: {t} - {e}")
        return 0  # return 0 if there's an error, or you could choose to handle it differently

def create_plot_and_slide(grouped, timestamp, prs, count, metric):
    # Create a figure and an axes.
    fig, ax = plt.subplots()

    for name, group in grouped:
        ax.plot(group['Hours Since Reported'], group[metric], label=name)

    # Setting the x-axis to show more recent times on the right
    ax.invert_xaxis()

    # Label the axes
    ax.set_xlabel('Time Passed (hours ago)')
    if metric == 'vBat':
        ylabel = 'Battery Voltage (vBat)'
    else:
        ylabel = 'RSRP'
    ax.set_ylabel(ylabel)

    # Title and legend
    ax.set_title(f'{metric} Over Time')
    ax.legend(title='IMEI')

    # Show a grid
    ax.grid(True)

    # Save the plot as an image
    image_dir = os.path.join(os.getcwd(), "images")
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)
    image_path = os.path.join(image_dir, f'{metric}_over_Time_{timestamp}_{count}.png')
    plt.savefig(image_path)
    plt.close(fig)

    # Add a new slide for the plot
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(0.1)
    slide.shapes.add_picture(image_path, left, top, width=Inches(10), height=Inches(8))

# Do statistical analysis for a given metric
def calculate_statistics(values):
    if values:
        count = len(values)
        mean = np.mean(values)
        std = np.std(values)
        min_val = min(values)
        max_val = max(values)
    else:
        count = mean = std = min_val = max_val = 'N/A'  # In case there are no valid values
    return count, mean, std, min_val, max_val

def create_table_slide(data, dev, prs, metric):
    values = [entry[metric] for entry in data]
    # Add a new slide for the summary table of statistics
    slide_layout = prs.slide_layouts[6]  # Choose a layout that fits a table well
    stats_slide = prs.slides.add_slide(slide_layout)

    # Define table dimensions
    rows, cols = 6, 2  # Additional row for headers
    left, top, width, height = Inches(3), Inches(3), Inches(5), Inches(0.2)  # Adjust as needed

    # Add a table to the slide
    table = stats_slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column headings
    table.cell(0, 0).text = 'IMEI'
    table.cell(1, 0).text = 'Count'
    table.cell(2, 0).text = 'Mean'
    table.cell(3, 0).text = 'Std'
    table.cell(4, 0).text = 'Min'
    table.cell(5, 0).text = 'Max'

    # Populate the table with data
    count, mean, std, min_val, max_val = calculate_statistics(values)
    table.cell(0, 1).text = dev
    table.cell(1, 1).text = str(count)
    table.cell(2, 1).text = f"{mean:.2f}" if mean != 'N/A' else 'N/A'
    table.cell(3, 1).text = f"{std:.2f}" if std != 'N/A' else 'N/A'
    table.cell(4, 1).text = str(min_val)
    table.cell(5, 1).text = str(max_val)

def to_excel(data_list, sheet_name, timestamp):
    if not data_list:  # Check if the data_list is empty
        print(f"No data to write for {sheet_name}")
        return
    
    df = pd.DataFrame(data_list)
    df = df[list(data_list[0].keys())]
    new_file_path = os.path.join(os.getcwd(), f'Upload Records Check {timestamp}.xlsx')
    if os.path.isfile(new_file_path) == False:
        df.to_excel(new_file_path, index=False, sheet_name=sheet_name)
    else:
        workbook = openpyxl.load_workbook(new_file_path)  # load workbook if already exists
        sheet = workbook.create_sheet(sheet_name)
        # append the dataframe results to the current excel file
        for row in dataframe_to_rows(df, header = True, index = False):
            sheet.append(row)
        workbook.save(new_file_path)  # save workbook
        workbook.close()  # close workbook

def run(fname):
    hours_ago = input("Enter the time period in hours (default 72): ")
    if hours_ago == "":
        hours_ago = 72

    # Read device list from file specified by the user
    with open(fname, 'r') as file:
        device_list = file.read().splitlines()
    print("reading device list: ", len(device_list))

    # An array tracking all the relevant data for all relevant devices
    rsrp_list = []
    vbat_list = []
    rsrp_group = []
    vbat_group = []
    slide_count = 0
    for i, dev in enumerate(device_list):
        print(f"---\nReport for device {dev}")
        data = get_device_data_v2(dev, int(hours_ago))
        if data is not None:
            for entry in data:
                try:
                    vbat_dict = vbat_data(entry, dev, hours_ago)
                    if vbat_dict:
                        vbat_group.append(vbat_dict)
                        vbat_list.append(vbat_dict)
                    rsrp_dict = rsrp_data(entry, dev, hours_ago)
                    if rsrp_dict:
                        rsrp_group.append(rsrp_dict)
                        rsrp_list.append(rsrp_dict)
                except:
                    print(f"Device {dev} shows error")
        df = pd.DataFrame(vbat_group)
        df['Hours Since Reported'] = df['Time passed since Reported'].apply(convert_to_hours)
        grouped = df.groupby('IMEI')
        create_plot_and_slide(grouped, datetime.now().strftime("%Y%m%d%H%M%S"), prs, slide_count, 'vBat')
        vbat_group = []
        slide_count += 1
        create_table_slide(rsrp_group, dev, prs, 'RSRP')
        slide_count += 1
        df2 = pd.DataFrame(rsrp_group)
        df2['Hours Since Reported'] = df2['Time passed since Reported'].apply(convert_to_hours)
        grouped2 = df2.groupby('IMEI')
        create_plot_and_slide(grouped2, datetime.now().strftime("%Y%m%d%H%M%S"), prs, slide_count, 'RSRP')
        rsrp_group = []
        slide_count += 1
    
    # store data in excel file
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    to_excel(vbat_list, "vBat", timestamp)
    to_excel(rsrp_list, "RSRP", timestamp)

    # Save the presentation
    prs.save(os.path.join(os.getcwd(), f'Presentation_{timestamp}.pptx'))
run(fname)